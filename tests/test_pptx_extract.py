from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

from opendocs._models import (
    HeadingBlock,
    InlineLink,
    InlineText,
    ListItemBlock,
    ListKind,
    ParagraphBlock,
    SpannedTableBlock,
    TableBlock,
)
from opendocs.parsers.office.models import ImageSlot, NativeSlot
from opendocs.parsers.office.pptx import extract_pptx
from opendocs.source import ParseWorkspace


def _save_presentation(presentation: Any, path: Path) -> Path:
    presentation.save(path)
    return path


def _paragraph_text(block: object) -> str:
    if not isinstance(block, ParagraphBlock | HeadingBlock | ListItemBlock):
        return ""
    return "".join(
        inline.text if isinstance(inline, InlineText) else inline.label for inline in block.inlines
    )


def _native_blocks(page) -> list[object]:
    return [block for slot in page.slots if isinstance(slot, NativeSlot) for block in slot.blocks]


def test_extract_pptx_preserves_slide_and_source_shape_order(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(7), Inches(5), Inches(2), Inches(1))
    first.text = "source first"
    second = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    second.text = "source second"
    presentation.slides.add_slide(presentation.slide_layouts[6])
    path = _save_presentation(presentation, tmp_path / "ordered.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))

    assert len(result.pages) == 2
    assert [_paragraph_text(block) for block in _native_blocks(result.pages[0])] == [
        "source first",
        "source second",
    ]
    assert result.pages[1].slots == ()
    assert [warning.code for warning in result.warnings] == ["pptx_blank_slide"]


def test_extract_pptx_recurses_groups_in_child_order(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    first = group.shapes.add_textbox(Inches(4), Inches(3), Inches(2), Inches(1))
    first.text = "group first"
    second = group.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    second.text = "group second"
    path = _save_presentation(presentation, tmp_path / "groups.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))

    assert [_paragraph_text(block) for block in _native_blocks(result.pages[0])] == [
        "group first",
        "group second",
    ]


def test_extract_pptx_preserves_placeholder_and_fallback_titles(tmp_path: Path) -> None:
    presentation = Presentation()
    titled = presentation.slides.add_slide(presentation.slide_layouts[0])
    titled.shapes.title.text = "Primary"
    titled.placeholders[1].text = "Subtitle"
    fallback = presentation.slides.add_slide(presentation.slide_layouts[6])
    numeric = fallback.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    numeric.text_frame.paragraphs[0].add_run().text = "2026"
    numeric.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    heading = fallback.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    heading.text_frame.paragraphs[0].add_run().text = "Fallback title"
    heading.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    body = fallback.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
    body.text_frame.paragraphs[0].add_run().text = "Body"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    path = _save_presentation(presentation, tmp_path / "titles.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))

    titled_blocks = _native_blocks(result.pages[0])
    assert [(type(block), _paragraph_text(block)) for block in titled_blocks] == [
        (HeadingBlock, "Primary"),
        (HeadingBlock, "Subtitle"),
    ]
    fallback_blocks = _native_blocks(result.pages[1])
    assert [
        (block.level if isinstance(block, HeadingBlock) else None, _paragraph_text(block))
        for block in fallback_blocks
    ] == [(None, "2026"), (1, "Fallback title"), (None, "Body")]


def test_extract_pptx_preserves_lists_soft_breaks_and_links(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    frame = shape.text_frame
    frame.clear()
    bullet = frame.paragraphs[0]
    bullet.text = "bullet"
    bullet.level = 1
    bullet._p.get_or_add_pPr().append(parse_xml(f'<a:buChar char="&#x2022;" {nsdecls("a")}/>'))
    ordered = frame.add_paragraph()
    ordered.text = "ordered"
    ordered._p.get_or_add_pPr().append(
        parse_xml(f'<a:buAutoNum type="arabicPeriod" startAt="3" {nsdecls("a")}/>')
    )
    linked = frame.add_paragraph()
    linked.add_run().text = "line one\vline two "
    safe = linked.add_run()
    safe.text = "safe"
    safe.hyperlink.address = "https://example.test"
    unsafe = linked.add_run()
    unsafe.text = "unsafe"
    unsafe.hyperlink.address = "javascript:alert(1)"
    path = _save_presentation(presentation, tmp_path / "text.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))
    blocks = _native_blocks(result.pages[0])

    assert isinstance(blocks[0], ListItemBlock)
    assert (blocks[0].kind, blocks[0].level, blocks[0].ordinal) == (ListKind.BULLET, 1, 1)
    assert isinstance(blocks[1], ListItemBlock)
    assert (blocks[1].kind, blocks[1].ordinal) == (ListKind.ORDERED, 3)
    assert isinstance(blocks[2], ParagraphBlock)
    assert _paragraph_text(blocks[2]) == "line one\nline two safeunsafe"
    assert [inline.target for inline in blocks[2].inlines if isinstance(inline, InlineLink)] == [
        "https://example.test",
        "javascript:alert(1)",
    ]


def test_extract_pptx_emits_simple_and_spanned_tables(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    simple_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(2))
    simple_shape.table.first_row = True
    simple_shape.table.cell(0, 0).text = "A"
    simple_shape.table.cell(0, 1).text = "B"
    simple_shape.table.cell(1, 0).text = "1"
    simple_shape.table.cell(1, 1).text = "2"
    merged_shape = slide.shapes.add_table(2, 2, Inches(5), Inches(1), Inches(3), Inches(2))
    merged_shape.table.cell(0, 0).merge(merged_shape.table.cell(0, 1))
    merged_shape.table.cell(0, 0).text = "wide"
    merged_shape.table.cell(1, 0).text = "left"
    merged_shape.table.cell(1, 1).text = "right"
    path = _save_presentation(presentation, tmp_path / "tables.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))
    blocks = _native_blocks(result.pages[0])

    assert blocks[0] == TableBlock((("A", "B"), ("1", "2")), 1)
    assert isinstance(blocks[1], SpannedTableBlock)
    assert blocks[1].cells[0].column_span == 2
    assert [cell.text for cell in blocks[1].cells] == ["wide", "left", "right"]


def test_extract_pptx_keeps_chart_title_and_data_adjacent(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["Q1", "Q2"]
    data.add_series("Revenue", (10, 20))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(3),
        data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly"
    path = _save_presentation(presentation, tmp_path / "chart.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))
    slot = result.pages[0].slots[0]

    assert isinstance(slot, NativeSlot)
    assert isinstance(slot.blocks[0], HeadingBlock)
    assert _paragraph_text(slot.blocks[0]) == "Quarterly"
    assert slot.blocks[1] == TableBlock(
        (("Series / Category", "Q1", "Q2"), ("Revenue", "10.0", "20.0")),
        1,
    )


def test_extract_pptx_emits_picture_at_shape_position(tmp_path: Path) -> None:
    image_path = tmp_path / "source.png"
    Image.new("RGB", (16, 8), "red").save(image_path, "PNG")
    image_bytes = image_path.read_bytes()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    before = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    before.text = "before"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(2), Inches(1))
    after = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
    after.text = "after"
    path = _save_presentation(presentation, tmp_path / "picture.pptx")
    workspace = ParseWorkspace(tmp_path / "workspace")

    result = extract_pptx(path, workspace)

    assert [type(slot) for slot in result.pages[0].slots] == [NativeSlot, ImageSlot, NativeSlot]
    image_slot = result.pages[0].slots[1]
    assert isinstance(image_slot, ImageSlot)
    assert image_slot.content_sha256 == hashlib.sha256(image_bytes).hexdigest()
    assert workspace.output_path(image_slot.artifact_name).read_bytes() == image_bytes
    assert image_slot.bbox.require_normalized() == image_slot.bbox


def test_extract_pptx_ignores_empty_decorative_shapes(tmp_path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    path = _save_presentation(presentation, tmp_path / "decorative.pptx")

    result = extract_pptx(path, ParseWorkspace(tmp_path / "workspace"))

    assert result.pages[0].slots == ()
    assert result.warnings == ()
