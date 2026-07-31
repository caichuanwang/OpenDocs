from __future__ import annotations

import asyncio
import io
import warnings
from collections.abc import Callable
from pathlib import Path
from typing import cast

import pytest
from docx import Document
from docx.enum.text import WD_BREAK
from pptx import Presentation
from pptx.util import Inches

import opendocs._runtime as runtime_module
from opendocs import (
    DocumentTimeoutError,
    LimitExceededError,
    OpenDocsWarning,
    ParseOptions,
    aparse,
    parse,
)
from opendocs.source import Source
from tests.native_worker_helpers import sleep_and_echo


class NamedBytesIO(io.BytesIO):
    def __init__(self, data: bytes, name: str) -> None:
        super().__init__(data)
        self.name = name


def _docx_bytes(*paragraphs: str, page_breaks: int = 0) -> bytes:
    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    for _ in range(page_breaks):
        document.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes(*slides: str) -> bytes:
    presentation = Presentation()
    for text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.text = text
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _source_factory(kind: str, path: Path, content: bytes) -> Callable[[], Source]:
    if kind == "path":
        path.write_bytes(content)
        return lambda: path
    if kind == "bytes":
        return lambda: content
    if kind == "named_stream":
        return lambda: NamedBytesIO(content, str(path))
    if kind == "unnamed_stream":
        return lambda: io.BytesIO(content)
    raise AssertionError(f"unknown source kind: {kind}")


@pytest.mark.parametrize("source_kind", ["path", "bytes", "named_stream", "unnamed_stream"])
@pytest.mark.parametrize(
    ("suffix", "content", "expected"),
    [
        (".docx", _docx_bytes("DOCX body"), "DOCX body\n"),
        (".pptx", _pptx_bytes("PPTX body"), "<!-- page: 1 -->\n\n# PPTX body\n"),
    ],
    ids=["docx", "pptx"],
)
def test_office_public_api_supports_all_input_shapes_with_sync_async_parity(
    tmp_path: Path,
    source_kind: str,
    suffix: str,
    content: bytes,
    expected: str,
) -> None:
    source = _source_factory(source_kind, tmp_path / f"source{suffix}", content)

    sync_input = source()
    async_input = source()
    assert parse(sync_input) == expected
    assert asyncio.run(aparse(async_input)) == expected
    if hasattr(sync_input, "closed"):
        assert sync_input.closed is False
    if hasattr(async_input, "closed"):
        assert async_input.closed is False


def test_pptx_max_pages_fails_before_rendering() -> None:
    with pytest.raises(LimitExceededError, match="1 page limit"):
        parse(_pptx_bytes("first", "second"), options=ParseOptions(max_pages=1))


def test_docx_max_pages_does_not_treat_hard_breaks_as_physical_pages() -> None:
    result = parse(_docx_bytes("before", "after", page_breaks=2), options=ParseOptions(max_pages=1))

    assert "before" in result
    assert "after" in result
    assert result.count("<!-- page-break -->") == 2


def test_office_output_limit_stops_between_complete_blocks() -> None:
    with warnings.catch_warnings(record=True) as captured:
        warnings.simplefilter("always", OpenDocsWarning)
        result = parse(
            _docx_bytes("first", "second block"),
            options=ParseOptions(max_output_chars=len("first\n")),
        )

    assert result == "first\n"
    assert [cast(OpenDocsWarning, warning.message).code for warning in captured] == [
        "output_truncated"
    ]


@pytest.mark.parametrize(
    ("content", "expected"),
    [
        (_docx_bytes("repeatable DOCX"), "repeatable DOCX\n"),
        (_pptx_bytes("repeatable PPTX"), "<!-- page: 1 -->\n\n# repeatable PPTX\n"),
    ],
    ids=["docx", "pptx"],
)
def test_office_public_api_is_deterministic_across_repeated_native_runs(
    content: bytes,
    expected: str,
) -> None:
    assert [parse(content) for _ in range(3)] == [expected, expected, expected]


@pytest.mark.asyncio
@pytest.mark.parametrize("source_kind", ["bytes", "named_stream", "unnamed_stream"])
async def test_office_external_cancellation_reaps_worker_and_cleans_owned_inputs(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    source_kind: str,
) -> None:
    original_run_native = runtime_module.ParserRuntime.run_native
    entered = asyncio.Event()
    source_path: Path | None = None
    workspace_path: Path | None = None
    runtime: runtime_module.ParserRuntime | None = None

    async def blocked_run_native(
        self: runtime_module.ParserRuntime,
        function: object,
        *args: object,
    ) -> object:
        nonlocal runtime, source_path, workspace_path
        del function
        runtime = self
        source_path = args[1] if isinstance(args[1], Path) else None
        workspace_path = args[2] if isinstance(args[2], Path) else None
        worker_task = asyncio.create_task(original_run_native(self, sleep_and_echo, 5.0, "done"))
        for _ in range(100):
            if self.native_worker.pid is not None:
                break
            await asyncio.sleep(0.01)
        entered.set()
        return await worker_task

    monkeypatch.setattr(runtime_module.ParserRuntime, "run_native", blocked_run_native)
    source = _source_factory(source_kind, tmp_path / "source.docx", _docx_bytes("body"))
    task = asyncio.create_task(aparse(source()))
    await asyncio.wait_for(entered.wait(), timeout=2)
    task.cancel()

    with pytest.raises(asyncio.CancelledError):
        await task

    assert runtime is not None
    assert runtime.native_worker.is_alive is False
    assert source_path is not None
    assert workspace_path is not None
    for _ in range(100):
        if not source_path.exists():
            break
        await asyncio.sleep(0.01)
    assert not source_path.exists()
    assert not workspace_path.exists()


@pytest.mark.asyncio
async def test_office_timeout_reaps_worker_and_cleans_owned_input(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    original_run_native = runtime_module.ParserRuntime.run_native
    source_path: Path | None = None
    workspace_path: Path | None = None
    runtime: runtime_module.ParserRuntime | None = None

    async def blocked_run_native(
        self: runtime_module.ParserRuntime,
        function: object,
        *args: object,
    ) -> object:
        nonlocal runtime, source_path, workspace_path
        del function
        runtime = self
        source_path = args[1] if isinstance(args[1], Path) else None
        workspace_path = args[2] if isinstance(args[2], Path) else None
        return await original_run_native(self, sleep_and_echo, 5.0, "done")

    monkeypatch.setattr(runtime_module.ParserRuntime, "run_native", blocked_run_native)

    with pytest.raises(DocumentTimeoutError):
        await aparse(_docx_bytes("body"), options=ParseOptions(timeout=0.1))

    assert runtime is not None
    assert runtime.native_worker.is_alive is False
    assert source_path is not None and not source_path.exists()
    assert workspace_path is not None and not workspace_path.exists()
