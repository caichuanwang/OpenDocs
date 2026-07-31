from __future__ import annotations

import hashlib
import re
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from opendocs._models import (
    BBox,
    Block,
    DocumentType,
    HeadingBlock,
    Inline,
    InlineLink,
    InlineText,
    ListItemBlock,
    ListKind,
    ParagraphBlock,
    SpannedTableBlock,
    SpannedTableCell,
    TableBlock,
    WarningRecord,
)
from opendocs.errors import CorruptDocumentError, RuntimeDependencyError
from opendocs.parsers.office.models import (
    ImageSlot,
    NativeSlot,
    OfficeDocument,
    OfficePage,
    document_to_wire,
)
from opendocs.parsers.office.package import (
    ExtractedMediaArtifact,
    extract_package_media,
    open_validated_office_document,
)
from opendocs.source import ParseWorkspace

_TITLE_PLACEHOLDERS = frozenset(
    {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.VERTICAL_TITLE}
)
_NUMERIC_TITLE = re.compile(r"[0-9\uFF10-\uFF19]+(?:[./:\uFF1A-][0-9\uFF10-\uFF19]+)*")
_UNSUPPORTED_SHAPES = frozenset(
    {
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.IGX_GRAPHIC,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
        MSO_SHAPE_TYPE.WEB_VIDEO,
    }
)


def _warning(code: str, slide_number: int, source_index: int) -> WarningRecord:
    return WarningRecord(
        code=code,
        message=f"PPTX slide {slide_number} source {source_index}: {code.replace('_', ' ')}",
    )


def _normalized_bbox(shape: Any, slide_width: int, slide_height: int) -> BBox:
    if slide_width <= 0 or slide_height <= 0:
        return BBox(0.0, 0.0, 1.0, 1.0)
    try:
        left = min(1.0, max(0.0, float(shape.left) / slide_width))
        top = min(1.0, max(0.0, float(shape.top) / slide_height))
        right = min(1.0, max(0.0, float(shape.left + shape.width) / slide_width))
        bottom = min(1.0, max(0.0, float(shape.top + shape.height) / slide_height))
    except (AttributeError, TypeError, ValueError):
        return BBox(0.0, 0.0, 1.0, 1.0)
    if left >= right or top >= bottom:
        return BBox(0.0, 0.0, 1.0, 1.0)
    return BBox(left, top, right, bottom)


def _append_text(inlines: list[Inline], text: str) -> None:
    normalized = (
        text.replace("_x000B_", "\n").replace("\v", "\n").replace("\r\n", "\n").replace("\r", "\n")
    )
    if not normalized:
        return
    if inlines and isinstance(inlines[-1], InlineText):
        previous = inlines[-1]
        inlines[-1] = InlineText(previous.text + normalized)
    else:
        inlines.append(InlineText(normalized))


def _paragraph_inlines(paragraph: Any) -> tuple[Inline, ...]:
    inlines: list[Inline] = []
    for run in paragraph.runs:
        text = (
            run.text.replace("_x000B_", "\n")
            .replace("\v", "\n")
            .replace("\r\n", "\n")
            .replace("\r", "\n")
        )
        if not text:
            continue
        target = run.hyperlink.address
        if target:
            try:
                inlines.append(InlineLink(text, target))
            except (TypeError, ValueError):
                _append_text(inlines, text)
        else:
            _append_text(inlines, text)
    if not inlines and paragraph.text:
        _append_text(inlines, paragraph.text)
    return tuple(inlines)


def _bullet_properties(paragraph: Any) -> tuple[ListKind, int] | None:
    properties = paragraph._p.pPr
    if properties is None:
        return None
    for child in properties:
        local_name = child.tag.rsplit("}", 1)[-1]
        if local_name == "buNone":
            return None
        if local_name == "buAutoNum":
            raw_start = child.get("startAt")
            start = int(raw_start) if raw_start and raw_start.isdigit() else 1
            return ListKind.ORDERED, max(start, 1)
        if local_name in {"buChar", "buBlip"}:
            return ListKind.BULLET, 1
    return None


def _text_blocks(shape: Any, heading_level: int | None, list_id: int) -> tuple[Block, ...]:
    blocks: list[Block] = []
    ordinals: dict[int, int] = {}
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        inlines = _paragraph_inlines(paragraph)
        if not inlines or not any(
            inline.text.strip() if isinstance(inline, InlineText) else inline.label.strip()
            for inline in inlines
        ):
            continue
        bullet = _bullet_properties(paragraph)
        if bullet is not None:
            kind, explicit_start = bullet
            level = max(0, int(paragraph.level))
            ordinal = explicit_start if kind is ListKind.ORDERED else 1
            if kind is ListKind.ORDERED and explicit_start == 1:
                ordinal = ordinals.get(level, 0) + 1
            ordinals[level] = ordinal
            blocks.append(ListItemBlock(list_id, level, kind, ordinal, inlines))
        elif heading_level is not None and paragraph_index == 0:
            blocks.append(HeadingBlock(heading_level, inlines))
        else:
            blocks.append(ParagraphBlock(inlines))
    return tuple(blocks)


def _iter_text_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_shapes(shape.shapes)
        elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            yield shape


def _shape_font_size(shape: Any) -> float | None:
    sizes: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size is not None:
            sizes.append(float(paragraph.font.size.pt))
        sizes.extend(float(run.font.size.pt) for run in paragraph.runs if run.font.size is not None)
    return max(sizes, default=None)


def _is_title_candidate(shape: Any) -> bool:
    text = "".join(shape.text_frame.text.split())
    return bool(text) and len(text) <= 160 and _NUMERIC_TITLE.fullmatch(text) is None


def _heading_levels(shapes: Iterable[Any]) -> dict[int, int]:
    text_shapes = list(_iter_text_shapes(shapes))
    levels: dict[int, int] = {}
    has_primary = False
    for shape in text_shapes:
        if not shape.is_placeholder:
            continue
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in _TITLE_PLACEHOLDERS:
            levels[shape.shape_id] = 1
            has_primary = True
        elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
            levels[shape.shape_id] = 2
    if has_primary:
        return levels

    candidates = [shape for shape in text_shapes if _is_title_candidate(shape)]
    ranked = [
        (size, -index, shape)
        for index, shape in enumerate(candidates)
        if (size := _shape_font_size(shape)) is not None
    ]
    if ranked:
        levels[max(ranked, key=lambda item: (item[0], item[1]))[2].shape_id] = 1
    elif len(candidates) == 1:
        levels[candidates[0].shape_id] = 1
    return levels


def _table_block(table: Any) -> TableBlock | SpannedTableBlock | None:
    row_count = len(table.rows)
    column_count = len(table.columns)
    if row_count == 0 or column_count == 0:
        return None
    has_spans = any(
        cell.is_spanned or (cell.is_merge_origin and (cell.span_height > 1 or cell.span_width > 1))
        for row in table.rows
        for cell in row.cells
    )
    header_rows = 1 if table.first_row else 0
    if not has_spans:
        return TableBlock(
            tuple(tuple(cell.text for cell in row.cells) for row in table.rows),
            header_rows,
        )

    cells: list[SpannedTableCell] = []
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            if cell.is_spanned:
                continue
            cells.append(
                SpannedTableCell(
                    row_index,
                    column_index,
                    int(cell.span_height) if cell.is_merge_origin else 1,
                    int(cell.span_width) if cell.is_merge_origin else 1,
                    cell.text,
                )
            )
    return SpannedTableBlock(row_count, column_count, tuple(cells), header_rows)


def _chart_blocks(shape: Any) -> tuple[Block, ...]:
    chart = shape.chart
    blocks: list[Block] = []
    if chart.has_title:
        title = chart.chart_title.text_frame.text.strip()
        if title:
            blocks.append(HeadingBlock(2, (InlineText(title),)))
    try:
        plot = chart.plots[0]
        categories = tuple(str(category) for category in plot.categories)
        series = tuple(chart.series)
        if categories and series:
            rows: list[tuple[str, ...]] = [("Series / Category", *categories)]
            for item in series:
                values = tuple("" if value is None else str(value) for value in item.values)
                width = len(categories)
                values = values[:width] + ("",) * max(0, width - len(values))
                rows.append((str(item.name or ""), *values))
            blocks.append(TableBlock(tuple(rows), 1))
    except (AttributeError, IndexError, KeyError, TypeError, ValueError):
        pass
    return tuple(blocks)


def _media_by_digest(
    artifacts: tuple[ExtractedMediaArtifact, ...],
) -> dict[str, ExtractedMediaArtifact]:
    result: dict[str, ExtractedMediaArtifact] = {}
    for artifact in artifacts:
        result.setdefault(artifact.content_sha256, artifact)
    return result


def _walk_shapes(
    shapes: Iterable[Any],
    *,
    slide_number: int,
    slide_width: int,
    slide_height: int,
    heading_levels: dict[int, int],
    media: dict[str, ExtractedMediaArtifact],
    slots: list[NativeSlot | ImageSlot],
    warnings: list[WarningRecord],
) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(
                shape.shapes,
                slide_number=slide_number,
                slide_width=slide_width,
                slide_height=slide_height,
                heading_levels=heading_levels,
                media=media,
                slots=slots,
                warnings=warnings,
            )
            continue

        source_index = len(slots)
        if getattr(shape, "has_table", False):
            table = _table_block(shape.table)
            if table is not None:
                slots.append(NativeSlot(source_index, (table,)))
            continue
        if getattr(shape, "has_chart", False):
            blocks = _chart_blocks(shape)
            if blocks:
                slots.append(NativeSlot(source_index, blocks))
            else:
                warnings.append(_warning("pptx_chart_data_unavailable", slide_number, source_index))
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                digest = hashlib.sha256(blob).hexdigest()
                artifact = media[digest]
                alt_text = shape.name or None
                slots.append(
                    ImageSlot(
                        source_index,
                        artifact.artifact_name,
                        digest,
                        _normalized_bbox(shape, slide_width, slide_height),
                        alt_text,
                    )
                )
            except (AttributeError, KeyError, OSError, TypeError, ValueError):
                warnings.append(_warning("pptx_image_unavailable", slide_number, source_index))
            continue
        if getattr(shape, "has_text_frame", False):
            blocks = _text_blocks(
                shape,
                heading_levels.get(shape.shape_id),
                list_id=(slide_number << 20) + source_index,
            )
            if blocks:
                slots.append(NativeSlot(source_index, blocks))
            continue
        if shape.shape_type in _UNSUPPORTED_SHAPES:
            warnings.append(_warning("pptx_unsupported_shape", slide_number, source_index))


def extract_pptx(path: Path, workspace: ParseWorkspace) -> OfficeDocument:
    if not isinstance(path, Path):
        raise TypeError("path must be a Path")
    if not isinstance(workspace, ParseWorkspace):
        raise TypeError("workspace must be a ParseWorkspace")
    try:
        artifacts = extract_package_media(
            path,
            document_type=DocumentType.PPTX,
            workspace=workspace,
        )
        presentation = open_validated_office_document(
            path,
            document_type=DocumentType.PPTX,
            opener=lambda source: Presentation(str(source)),
        )
    except ImportError as error:
        raise RuntimeDependencyError("python-pptx is required to parse PPTX documents") from error
    except (CorruptDocumentError, RuntimeDependencyError):
        raise
    except (KeyError, OSError, TypeError, ValueError) as error:
        raise CorruptDocumentError("PPTX package could not be opened") from error

    media = _media_by_digest(artifacts)
    slide_width = int(presentation.slide_width or 0)
    slide_height = int(presentation.slide_height or 0)
    pages: list[OfficePage] = []
    warnings: list[WarningRecord] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slots: list[NativeSlot | ImageSlot] = []
        _walk_shapes(
            slide.shapes,
            slide_number=slide_number,
            slide_width=slide_width,
            slide_height=slide_height,
            heading_levels=_heading_levels(slide.shapes),
            media=media,
            slots=slots,
            warnings=warnings,
        )
        pages.append(OfficePage(slide_number, tuple(slots)))

    if any(page.slots for page in pages):
        for page in pages:
            if not page.slots:
                warnings.append(_warning("pptx_blank_slide", page.page_number, 0))
    return OfficeDocument(DocumentType.PPTX, tuple(pages), tuple(warnings))


def extract_pptx_to_wire(path: Path, workspace_path: Path) -> dict[str, object]:
    return document_to_wire(extract_pptx(path, ParseWorkspace(workspace_path)))
