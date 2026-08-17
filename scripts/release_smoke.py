from __future__ import annotations

import argparse
import asyncio
import io
import warnings
from datetime import date
from importlib.metadata import version
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from opendocs import OpenDocsWarning, aparse, parse


def run_smoke(directory: Path, *, expected_version: str) -> dict[str, bool]:
    directory.mkdir(parents=True, exist_ok=True)
    if version("opendocs-sdk") != expected_version:
        raise RuntimeError(f"installed opendocs-sdk version is not {expected_version}")

    text_path = directory / "smoke.txt"
    text_path.write_text("OpenDocs text smoke", encoding="utf-8")
    markdown_path = directory / "smoke.md"
    markdown_path.write_text("# OpenDocs Markdown smoke\n", encoding="utf-8")
    pdf_path = directory / "smoke.pdf"
    pdf_path.write_bytes(
        _native_pdf("OpenDocs native PDF smoke text with enough words for reliable extraction")
    )
    docx_path = directory / "smoke.docx"
    document = Document()
    document.add_heading("OpenDocs DOCX smoke")
    document.add_paragraph("Native Office extraction")
    document.save(str(docx_path))
    pptx_path = directory / "smoke.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text = "OpenDocs PPTX smoke"
    presentation.save(str(pptx_path))
    xlsx_path = directory / "smoke.xlsx"
    _write_xlsx_smoke(xlsx_path)

    pdf_markdown = " ".join(parse(pdf_path).split())
    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always", OpenDocsWarning)
        xlsx_markdown = parse(xlsx_path)
    xlsx_warning_codes = [
        warning.message.code for warning in caught if isinstance(warning.message, OpenDocsWarning)
    ]
    results = {
        "text": "OpenDocs text smoke" in parse(text_path),
        "markdown": "# OpenDocs Markdown smoke" in parse(markdown_path),
        "pdf": "OpenDocs native PDF smoke" in pdf_markdown,
        "docx": "OpenDocs DOCX smoke" in parse(docx_path),
        "pptx": "OpenDocs PPTX smoke" in parse(pptx_path),
        "async": "OpenDocs text smoke" in asyncio.run(aparse(text_path)),
        "xlsx": all(
            anchor in xlsx_markdown
            for anchor in (
                "# Summary (Visible)",
                "# Hidden (Hidden)",
                "# Very Hidden (Very Hidden)",
                "# Empty (Visible)",
                "$1,234.50",
                "2026-08-14",
                '<td colspan="2">Merged smoke</td>',
                "=SUM(B1,1)",
            )
        )
        and xlsx_warning_codes == ["xlsx_formula_cache_missing"],
    }
    if not all(results.values()):
        failed = sorted(name for name, passed in results.items() if not passed)
        raise RuntimeError(f"release smoke failed for: {', '.join(failed)}")
    return dict(sorted(results.items()))


def _write_xlsx_smoke(path: Path) -> None:
    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary["A1"] = "Amount"
    summary["B1"] = 1234.5
    summary["B1"].number_format = "$#,##0.00"
    summary["A2"] = "Date"
    summary["B2"] = date(2026, 8, 14)
    summary["B2"].number_format = "yyyy-mm-dd"
    summary.merge_cells("A3:B3")
    summary["A3"] = "Merged smoke"
    summary["A4"] = "=SUM(B1,1)"

    hidden = workbook.create_sheet("Hidden")
    hidden.sheet_state = "hidden"
    hidden["A1"] = "Hidden smoke"
    very_hidden = workbook.create_sheet("Very Hidden")
    very_hidden.sheet_state = "veryHidden"
    very_hidden["A1"] = "Very hidden smoke"
    workbook.create_sheet("Empty")
    workbook.save(path)
    workbook.close()

    rewritten = path.with_suffix(".rewritten.xlsx")
    formula_with_empty_cache = b"<f>SUM(B1,1)</f><v></v>"
    replacement_count = 0
    with ZipFile(path) as source, ZipFile(rewritten, "w", ZIP_DEFLATED) as target:
        for member in source.infolist():
            payload = source.read(member.filename)
            if member.filename == "xl/worksheets/sheet1.xml":
                replacement_count = payload.count(formula_with_empty_cache)
                payload = payload.replace(formula_with_empty_cache, b"<f>SUM(B1,1)</f>")
            target.writestr(member, payload)
    if replacement_count != 1:
        raise RuntimeError("XLSX smoke fixture did not contain the expected empty formula cache")
    rewritten.replace(path)


def _native_pdf(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 14 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode() + stream + b"\nendstream",
    )
    output = io.BytesIO()
    output.write(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(output.tell())
        output.write(f"{number} 0 obj\n".encode())
        output.write(body)
        output.write(b"\nendobj\n")
    xref_offset = output.tell()
    output.write(f"xref\n0 {len(objects) + 1}\n".encode())
    output.write(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.write(f"{offset:010d} 00000 n \n".encode())
    output.write(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return output.getvalue()


def _parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Run native-only smoke checks against an installed OpenDocs artifact."
    )
    parser.add_argument("directory", type=Path)
    parser.add_argument("--version")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = _parser().parse_args(argv)
    results = run_smoke(args.directory, expected_version=args.version or version("opendocs-sdk"))
    print(", ".join(f"{name}=pass" for name in results))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
